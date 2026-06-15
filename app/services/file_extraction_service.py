from __future__ import annotations


def detect_file_type(file_name: str) -> str | None:
    lowered = file_name.lower()
    if lowered.endswith('.pdf'):
        return 'pdf'
    if lowered.endswith('.pptx'):
        return 'pptx'
    if lowered.endswith('.docx'):
        return 'docx'
    return None


def extract_text_from_pdf(file_path: str) -> str:
    import fitz
    parts: list[str] = []
    with fitz.open(file_path) as doc:
        for page in doc:
            parts.append(page.get_text())
    return chr(10).join(parts)


def extract_text_from_pptx(file_path: str) -> str:
    from pptx import Presentation
    presentation = Presentation(file_path)
    parts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    return chr(10).join(parts)


def extract_text_from_docx(file_path: str) -> str:
    from docx import Document
    document = Document(file_path)
    parts = [paragraph.text for paragraph in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                parts.append(cell.text)
    return chr(10).join(parts)


def extract_text_by_file_type(file_path: str, file_type: str) -> str:
    if file_type == 'pdf':
        return extract_text_from_pdf(file_path)
    if file_type == 'pptx':
        return extract_text_from_pptx(file_path)
    if file_type == 'docx':
        return extract_text_from_docx(file_path)
    raise ValueError(f'Unsupported file type: {file_type}')


def extract_text(file_path: str, file_type: str) -> str:
    return extract_text_by_file_type(file_path, file_type)
